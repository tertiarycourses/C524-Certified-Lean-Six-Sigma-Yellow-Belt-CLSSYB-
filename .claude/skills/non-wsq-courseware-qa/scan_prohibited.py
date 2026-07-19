#!/usr/bin/env python3
"""Scan non-WSQ courseware for prohibited WSQ/funding/assessment content.

Non-WSQ courses carry NO assessment and NO funding-compliance material. This
scanner reads the generated PPTX/DOCX/MD artifacts and reports every hit with
its exact location (slide number / paragraph index / line number), so a leak
can be traced straight back to the source.

Usage:
    python3 scan_prohibited.py <course-repo>      # defaults to cwd
Exit status: 0 = clean, 1 = prohibited content found.
"""
import os, re, sys, glob

# (label, regex) — word-boundary anchored to avoid false hits such as
# "assessment" inside "reassessment of the process" in genuine course content.
PATTERNS = [
    # (?<!non-) so a deliberate disclaimer like "Non-WSQ" is not a hit.
    ("WSQ",                 r"(?<!non-)(?<!non )\bWSQ\b"),
    ("SSG",                 r"\bSSG\b"),
    ("SkillsFuture",        r"\bSkills\s*Future\b"),
    ("TRAQOM",              r"\bTRAQOM\b"),
    ("TGS course ref",      r"\bTGS-\d+"),
    ("digital attendance",  r"\bdigital attendance\b"),
    # Course-funding language only — a course may legitimately discuss project
    # funding as subject matter, so require the training-funding context.
    ("course funding/subsidy", r"\b(course|training|SSG|WSQ|government|nett?)[- ]?(fee|funding|funded|subsidy|subsidised|subsidized)\b"
                               r"|\b(funding|subsidy)\s+(eligib|support|scheme|amount|rate)"
                               r"|\bfunded\s+(course|programme|program|training)\b"
                               r"|\bSkillsFuture\s+Credit\b"),
    ("75% attendance",      r"\b75\s*%\s*attendance\b"),
    ("Written Assessment",  r"\bWritten Assessment\b|\bWA\s*\(SAQ\)|\bSAQ\b"),
    ("Practical Performance", r"\bPractical Performance\b|\bPP\s*\("),
    # NOTE: a bare \bassessment\b is deliberately NOT flagged — many courses
    # legitimately teach "risk assessment", "impact assessment", "FMEA
    # assessment" etc. Only FORMAL-assessment phrasing is prohibited.
    ("formal assessment",   r"\b(final|formal|summative|course|open[- ]book)\s+assessment\b"
                            r"|\bassessment\s+(flow|day|paper|script|brief(ing)?|criteria|record|plan)\b"
                            r"|\bbriefing\s+for\s+assessment\b"
                            r"|\beligible\s+for\s+assessment\b"),
    ("assessor/marking",    r"\bassessor\b|\bmarking guide\b|\banswer key\b|\bmodel answer\b"),
    ("Assessment Summary Record", r"\bAssessment Summary Record\b"),
    ("competency verdict",  r"\b(competent|not yet competent|NYC)\b"),
]
COMPILED = [(lab, re.compile(rx, re.I)) for lab, rx in PATTERNS]


def _hits(text):
    return [lab for lab, rx in COMPILED if rx.search(text or "")]


def scan_pptx(path):
    from pptx import Presentation
    out = []
    prs = Presentation(path)
    for i, slide in enumerate(prs.slides, 1):
        parts = []
        for sh in slide.shapes:
            if sh.has_text_frame:
                parts.append(sh.text_frame.text)
            if getattr(sh, "has_table", False) and sh.has_table:
                for row in sh.table.rows:
                    parts += [c.text for c in row.cells]
        text = "\n".join(parts)
        for lab in _hits(text):
            out.append((f"slide {i}", lab, _excerpt(text, lab)))
    return out


def scan_docx(path):
    from docx import Document
    out = []
    doc = Document(path)
    for i, p in enumerate(doc.paragraphs, 1):
        for lab in _hits(p.text):
            out.append((f"para {i}", lab, _excerpt(p.text, lab)))
    for ti, t in enumerate(doc.tables, 1):
        for ri, row in enumerate(t.rows, 1):
            text = " | ".join(c.text for c in row.cells)
            for lab in _hits(text):
                out.append((f"table {ti} row {ri}", lab, _excerpt(text, lab)))
    return out


def scan_text(path):
    out = []
    with open(path, encoding="utf-8", errors="replace") as fh:
        for i, line in enumerate(fh, 1):
            for lab in _hits(line):
                out.append((f"line {i}", lab, line.strip()[:120]))
    return out


def _excerpt(text, label):
    rx = dict((l, r) for l, r in PATTERNS).get(label)
    m = re.search(rx, text, re.I) if rx else None
    if not m:
        return text.strip()[:120]
    a, b = max(0, m.start() - 50), min(len(text), m.end() + 50)
    return text[a:b].replace("\n", " ").strip()[:140]


def main():
    repo = sys.argv[1] if len(sys.argv) > 1 else os.getcwd()
    targets = []
    for pat in ("courseware/*.pptx", "courseware/*.docx", "*.md",
                "labs/*.md", "labs/**/*.md"):
        targets += glob.glob(os.path.join(repo, pat), recursive=True)
    # Repo documentation is not learner-facing courseware. README/CLAUDE files
    # legitimately name the prohibited terms in order to state they are excluded.
    targets = [t for t in targets
               if os.path.basename(t).lower() not in
               ("readme.md", "claude.md", "contributing.md", "changelog.md")
               or os.path.basename(os.path.dirname(t)) == "labs"]
    # Skip superseded versions, Office lock files, and reference/ — the latter
    # holds WSQ source documents (e.g. the original Course Proposal) that are
    # inputs to the course, never learner-facing courseware.
    targets = [t for t in sorted(set(targets))
               if "/archive/" not in t and "/reference/" not in t and "~$" not in t]

    if not targets:
        print(f"No courseware artifacts found under {repo}")
        return 1

    total = 0
    for path in targets:
        ext = os.path.splitext(path)[1].lower()
        try:
            hits = (scan_pptx(path) if ext == ".pptx" else
                    scan_docx(path) if ext == ".docx" else
                    scan_text(path))
        except Exception as e:                      # keep scanning the rest
            print(f"!! could not read {os.path.relpath(path, repo)}: {e}")
            continue
        if hits:
            print(f"\n=== {os.path.relpath(path, repo)} — {len(hits)} hit(s)")
            for loc, lab, ex in hits:
                print(f"  [{lab}] {loc}: {ex}")
            total += len(hits)

    print("\n" + "=" * 60)
    if total:
        print(f"FAIL — {total} prohibited-content hit(s) across {len(targets)} artifact(s).")
        print("Non-WSQ courseware must contain NO assessment, funding, SSG or TRAQOM content.")
        return 1
    print(f"PASS — {len(targets)} artifact(s) scanned, no prohibited content found.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
