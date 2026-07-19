#!/usr/bin/env python3
"""Generate a NON-WSQ course slide deck (all-white Tertiary house style).

Design helpers match the Tertiary house system (cover, section, content,
two_col, cards3, big_statement, step_slide, test_slide, brk). Content is driven
entirely by course_data.py + data_domainN.py so the deck stays 100% aligned with
the LP, LG and labs.

NON-WSQ: contains NO assessment, funding, SSG/SkillsFuture, TRAQOM or digital
attendance content. The assessment block a WSQ deck would carry is replaced by a
"How You'll Learn" flow, and the deck closes on the wrap-up + Thank You.
"""
import os, sys, re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
import course_data as C

# Import data_domain1..N dynamically so a course can have ANY number of domains
# (the WSQ engine hard-coded exactly five). Domains are loaded in order and
# concatenated into one contiguous activity list.
import importlib, glob as _glob
def _load_domains():
    acts=[]
    files=sorted(_glob.glob(os.path.join(HERE,"data_domain[0-9]*.py")),
                 key=lambda p:int("".join(ch for ch in os.path.basename(p) if ch.isdigit()) or 0))
    for f in files:
        mod=importlib.import_module(os.path.basename(f)[:-3])
        n="".join(ch for ch in os.path.basename(f) if ch.isdigit())
        acts+=getattr(mod,f"DOMAIN{n}",[])
    return acts
ACTIVITIES = _load_domains()

def _find_repo(start):
    """Locate the course repo (a dir containing both courseware/ and labs/).
    Env COURSE_REPO overrides. Keeps the build working wherever the skill lives."""
    env = os.environ.get("COURSE_REPO")
    if env and os.path.isdir(env):
        return env
    d = start
    for _ in range(8):
        d = os.path.dirname(d)
        if os.path.isdir(os.path.join(d, "courseware")) and os.path.isdir(os.path.join(d, "labs")):
            return d
    return os.path.dirname(os.path.dirname(HERE))
REPO = _find_repo(HERE)
ASSETS = os.path.join(os.path.dirname(HERE), "assets")   # co-located with the skill

# ---------------- palette (matches reference) ----------------
BLUE=RGBColor(0x1F,0x6F,0xEB); TEAL=RGBColor(0x10,0xB9,0x81); AMBER=RGBColor(0xF5,0x9E,0x0B)
INK=RGBColor(0x16,0x1B,0x26); GREY=RGBColor(0x5B,0x63,0x72); LIGHT=RGBColor(0xF5,0xF8,0xFC)
WHITE=RGBColor(0xFF,0xFF,0xFF); LINE=RGBColor(0xE2,0xE8,0xF0); VIOLET=RGBColor(0x7C,0x3A,0xED)
ORANGE=RGBColor(0xEA,0x58,0x0C); RED=RGBColor(0xDC,0x26,0x26)
# Per-topic accent themes for the section covers — cycled so ANY number of
# topics gets a distinct colour without hard-coding a specific course.
_ACCENTS=[ORANGE,RED,BLUE,TEAL,VIOLET,AMBER]
CARD_COLORS=[BLUE,TEAL,VIOLET]
TOPIC_THEME={t["num"]:_ACCENTS[i%len(_ACCENTS)] for i,t in enumerate(C.TOPICS)}
# Topic → its labs. Hoisted here (the WSQ engine defined it late) so the admin
# overview slides can render the schedule straight from the data.
TOPIC_ACTS = {t["num"]: [a for a in ACTIVITIES if a["topic"]==t["num"]] for t in C.TOPICS}

# ---------------- dark "masterclass" palette (opt-in via course_data.DARK_THEME) ----------------
DK_BG   = RGBColor(0x0D,0x0E,0x0B)   # slide background
DK_PANEL= RGBColor(0x15,0x17,0x12)   # cards/rows/panels
DK_AMBER= RGBColor(0xF5,0xA6,0x23)   # headlines, kickers, accents, number chips
DK_IVORY= RGBColor(0xE8,0xE6,0xDD)   # body text
DK_FAINT= RGBColor(0x2A,0x2C,0x28)   # divider rules
DK_DIM  = RGBColor(0x8A,0x8D,0x86)   # captions/footers
DK_GREEN= RGBColor(0x4A,0xDE,0x80)   # success/test-it accents

# Module-level theme switch — dark "masterclass" theme for the ENTIRE deck.
THEME={"dark":False}
def _bg():   return DK_BG    if THEME["dark"] else WHITE
def _panel():return DK_PANEL if THEME["dark"] else LIGHT
def _ink():  return DK_IVORY if THEME["dark"] else INK
def _grey(): return DK_DIM   if THEME["dark"] else GREY
def _line(): return DK_FAINT if THEME["dark"] else LINE
def _acc(c): return DK_AMBER if THEME["dark"] else c
# Non-WSQ house standard is the ALL-WHITE deck. A course may opt into the dark
# "masterclass" look with course_data.DARK_THEME = True.
THEME["dark"]=bool(getattr(C,"DARK_THEME",False))

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height
BLANK=prs.slide_layouts[6]

def slide(): return prs.slides.add_slide(BLANK)
def rect(s,x,y,w,h,color,line=None):
    sp=s.shapes.add_shape(1,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1)
    sp.shadow.inherit=False; return sp
def oval(s,x,y,w,h,color):
    sp=s.shapes.add_shape(9,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color
    sp.line.fill.background(); sp.shadow.inherit=False; return sp
def txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,space=4):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,line in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(space)
        for t,sz,col,bold in line:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=bold
            r.font.color.rgb=col; r.font.name="Arial"
    return tb
def bullets(s,x,y,w,h,items,size=18,color=None,gap=10,mcolor=BLUE):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    body=color if color is not None else _ink()
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(gap)
        lvl=it[1] if isinstance(it,tuple) else 0
        text=it[0] if isinstance(it,tuple) else it
        r=p.add_run(); r.text=("•  " if lvl==0 else "–  ")+text
        r.font.size=Pt(size if lvl==0 else size-2); r.font.color.rgb=body if lvl==0 else _grey()
        r.font.name="Arial"; r.font.bold=(lvl==0 and isinstance(it,tuple) and len(it)>2 and it[2])
    return tb

PAGE={"n":0}
def footer(s):
    PAGE["n"]+=1
    txt(s,Inches(0.4),Inches(7.05),Inches(7.5),Inches(0.35),
        [[(f"{C.SHORT_TITLE}  ·  {C.COURSE_CODE}",9,_grey(),False)]])
    txt(s,Inches(5.0),Inches(7.05),Inches(3.3),Inches(0.35),
        [[("© 2026 Tertiary Infotech Academy Pte Ltd",9,_grey(),False)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(12.4),Inches(7.05),Inches(0.6),Inches(0.35),
        [[(str(PAGE["n"]),9,_grey(),False)]],align=PP_ALIGN.RIGHT)
def head(s,title,kicker=None,kcolor=BLUE):
    rect(s,0,0,SW,SH,_bg()); rect(s,0,0,Inches(0.28),Inches(1.55),_acc(kcolor))
    if kicker: txt(s,Inches(0.85),Inches(0.48),Inches(11.9),Inches(0.4),[[(kicker,14,_acc(kcolor),True)]])
    # Auto-fit the title so long lab titles never wrap into the divider rule below.
    L=len(str(title)); tsize = 29 if L<=42 else (25 if L<=56 else 21)
    txt(s,Inches(0.85),Inches(0.9),Inches(11.95),Inches(0.82),[[(title,tsize,_ink(),True)]],anchor=MSO_ANCHOR.MIDDLE)
    rect(s,Inches(0.85),Inches(1.8),Inches(11.63),Inches(0.02),_line())
    return s
def _logo(name):
    p=os.path.join(ASSETS,name)
    return p if os.path.exists(p) else None

# ---------------- slide templates ----------------
def cover():
    s=slide(); rect(s,0,0,SW,SH,_bg())
    rect(s,0,0,SW,Inches(0.22),BLUE); rect(s,0,Inches(7.28),SW,Inches(0.22),TEAL)
    org=_logo("tertiary-infotech-logo.png")
    if org: s.shapes.add_picture(org,Inches(0.85),Inches(0.7),height=Inches(1.05))
    # course badge (top-right) — optional course badge asset; never a WSQ badge
    badge=_logo("course-badge.png")
    if badge:
        s.shapes.add_picture(badge,Inches(10.05),Inches(0.6),width=Inches(2.5))
    txt(s,Inches(0.9),Inches(2.3),Inches(12),Inches(0.6),[[("COURSE SLIDES",16,_acc(BLUE),True)]])
    txt(s,Inches(0.9),Inches(2.85),Inches(12.0),Inches(1.9),[[(C.TITLE,40,_ink(),True)]])
    rect(s,Inches(0.92),Inches(4.75),Inches(2.4),Inches(0.06),TEAL)
    txt(s,Inches(0.9),Inches(5.05),Inches(12),Inches(1.4),
        [[(f"Course Code: {C.COURSE_CODE}",16,_grey(),False)],
         [("Conducted by Tertiary Infotech Academy Pte Ltd  ·  UEN 201200696W",14,_grey(),False)]],space=6)
    txt(s,Inches(0.9),Inches(6.5),Inches(12),Inches(0.4),[[(f"Version {C.VERSION}  ·  {C.VERSION_DATE}",12,_grey(),False)]])
    txt(s,Inches(0.9),Inches(6.85),Inches(12),Inches(0.34),[[("© 2026 Tertiary Infotech Academy Pte Ltd. All rights reserved.  ·  www.tertiarycourses.com.sg",10,_grey(),False)]])

def section(kicker,title,n,sub="",accent=BLUE):
    s=slide(); rect(s,0,0,SW,SH,_bg()); rect(s,0,0,Inches(0.28),SH,accent)
    rect(s,Inches(0.85),Inches(2.5),Inches(0.14),Inches(2.0),accent)
    txt(s,Inches(1.25),Inches(2.55),Inches(11),Inches(0.6),[[(kicker,18,accent,True)]])
    txt(s,Inches(1.25),Inches(3.0),Inches(11.4),Inches(1.6),[[(title,40,_ink(),True)]])
    if sub: txt(s,Inches(1.27),Inches(4.55),Inches(11),Inches(0.8),[[(sub,16,_grey(),False)]])
    ghost=DK_FAINT if THEME["dark"] else RGBColor(0xE2,0xE8,0xF0)
    txt(s,Inches(10.0),Inches(0.7),Inches(2.8),Inches(1.6),[[(n,72,ghost,True)]],align=PP_ALIGN.RIGHT)
    footer(s)
def content(title,items,kicker=None,size=20):
    # HARD RULE: no bullet-wall slides — every content slide renders as numbered
    # visual cards (LIGHT card, accent bar, number chip), never plain bullets.
    import math as _m
    s=head(slide(),title,kicker)
    xs=[(i[0] if isinstance(i,tuple) else i) for i in items]
    n=max(1,len(xs)); cols=2 if n>3 else 1; rows=_m.ceil(n/cols)
    x0,y0=Inches(0.85),Inches(2.0); gw=Inches(11.63); gap=Inches(0.18)
    cw=int((gw-gap*(cols-1))/cols); ch=int((Inches(4.7)-gap*(rows-1))/rows)
    fs=(16 if rows<=2 else 14 if rows==3 else 12) if cols==2 else (17 if rows<=3 else 14)
    for idx,it in enumerate(xs):
        r,c=divmod(idx,cols)
        cx=x0+c*(cw+gap); cy=y0+r*(ch+gap)
        rect(s,cx,cy,cw,ch,_panel()); rect(s,cx,cy,Inches(0.07),ch,_acc(BLUE))
        oval(s,cx+Inches(0.2),int(cy+ch/2-Inches(0.19)),Inches(0.38),Inches(0.38),_acc(BLUE))
        txt(s,cx+Inches(0.2),int(cy+ch/2-Inches(0.17)),Inches(0.38),Inches(0.34),
            [[(str(idx+1),13,DK_BG if THEME["dark"] else WHITE,True)]],align=PP_ALIGN.CENTER)
        txt(s,cx+Inches(0.74),cy+Inches(0.08),cw-Inches(0.98),ch-Inches(0.16),
            [[(it,fs,_ink(),False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def two_col(title,left,right,kicker=None,lhead="",rhead=""):
    s=head(slide(),title,kicker)
    rect(s,Inches(0.85),Inches(1.95),Inches(5.7),Inches(4.7),_panel()); rect(s,Inches(6.95),Inches(1.95),Inches(5.55),Inches(4.7),_panel())
    if lhead: txt(s,Inches(1.1),Inches(2.15),Inches(5.2),Inches(0.4),[[(lhead,16,_acc(BLUE),True)]])
    if rhead: txt(s,Inches(7.2),Inches(2.15),Inches(5.0),Inches(0.4),[[(rhead,16,_acc(TEAL),True)]])
    bullets(s,Inches(1.1),Inches(2.7),Inches(5.2),Inches(3.8),left,size=16,color=_ink())
    bullets(s,Inches(7.2),Inches(2.7),Inches(5.05),Inches(3.8),right,size=16,color=_ink(),mcolor=TEAL); footer(s); return s
def cards3(title,cards,kicker):
    s=head(slide(),title,kicker); xs=[Inches(0.85),Inches(5.0),Inches(9.15)]
    for i,c in enumerate(cards[:3]):
        x=xs[i]; col=c[0]
        rect(s,x,Inches(1.95),Inches(3.65),Inches(4.7),_panel()); rect(s,x,Inches(1.95),Inches(3.65),Inches(0.12),col)
        txt(s,x+Inches(0.25),Inches(2.2),Inches(3.2),Inches(0.6),[[(c[1],19,col,True)]])
        bullets(s,x+Inches(0.25),Inches(2.95),Inches(3.2),Inches(3.4),c[2],size=14,color=_ink(),mcolor=col,gap=9)
    footer(s); return s
def big_statement(line1,line2,kicker,color=BLUE):
    s=slide(); rect(s,0,0,SW,SH,_bg()); rect(s,0,0,Inches(0.28),SH,color)
    txt(s,Inches(1.1),Inches(2.2),Inches(11),Inches(0.5),[[(kicker,16,color,True)]])
    txt(s,Inches(1.1),Inches(2.8),Inches(11.3),Inches(2.4),[[(line1,38,_ink(),True)]])
    if line2: txt(s,Inches(1.12),Inches(4.9),Inches(11),Inches(1.2),[[(line2,20,_grey(),False)]])
    footer(s); return s
import math
PALETTE=[BLUE,TEAL,VIOLET,AMBER]
def tile_grid(title,items,kicker=None,cols=2,size=15,icons=None,accent=BLUE):
    """Grid of light panels, each with a coloured icon/number badge + text.
    items: list of strings (or (title,caption) tuples). Much richer than a bullet list."""
    s=head(slide(),title,kicker,kcolor=accent)
    n=len(items); rows=math.ceil(n/cols)
    X0=Inches(0.85); Y0=Inches(1.95); TOTW=Inches(11.63); AREAH=Inches(4.78)
    gx=Inches(0.3); gy=Inches(0.26)
    cw=int((TOTW-gx*(cols-1))/cols); ch=int((AREAH-gy*(rows-1))/rows)
    bd=Inches(0.6)
    for i,it in enumerate(items):
        r=i//cols; c=i%cols
        x=int(X0+(cw+gx)*c); y=int(Y0+(ch+gy)*r); col=PALETTE[i%len(PALETTE)]
        rect(s,x,y,cw,ch,_panel()); rect(s,x,y,Inches(0.1),ch,col)
        oval(s,x+Inches(0.28),int(y+ch/2-bd/2),bd,bd,col)
        ic=icons[i] if icons else str(i+1)
        txt(s,x+Inches(0.28),int(y+ch/2-bd/2),bd,bd,[[(ic,19,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        tx=x+Inches(1.08); tw=cw-Inches(1.32)
        if isinstance(it,tuple):
            txt(s,tx,int(y+Inches(0.14)),tw,int(ch-Inches(0.2)),
                [[(it[0],size+2,_ink(),True)],[(it[1],size-2,_grey(),False)]],anchor=MSO_ANCHOR.MIDDLE,space=3)
        else:
            txt(s,tx,int(y+Inches(0.1)),tw,int(ch-Inches(0.16)),[[(it,size,_ink(),False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def flow_h(title,steps,kicker=None,color=BLUE):
    """Horizontal numbered flow: coloured chips connected by chevrons."""
    s=head(slide(),title,kicker,kcolor=color)
    n=len(steps); X0=Inches(0.85); TOTW=Inches(11.63); gap=Inches(0.34)
    cw=int((TOTW-gap*(n-1))/n); y=Inches(2.55); ch=Inches(3.15); bd=Inches(0.82)
    for i,st in enumerate(steps):
        x=int(X0+(cw+gap)*i)
        rect(s,x,y,cw,ch,_panel()); rect(s,x,y,cw,Inches(0.1),color)
        oval(s,int(x+cw/2-bd/2),int(y+Inches(0.42)),bd,bd,color)
        txt(s,int(x+cw/2-bd/2),int(y+Inches(0.42)),bd,bd,[[(str(i+1),30,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        txt(s,x+Inches(0.16),int(y+Inches(1.55)),cw-Inches(0.32),int(ch-Inches(1.7)),[[(st,14,_ink(),False)]],align=PP_ALIGN.CENTER)
        if i<n-1:
            txt(s,int(x+cw-Inches(0.04)),int(y+ch/2-Inches(0.3)),int(gap+Inches(0.08)),Inches(0.6),
                [[("▶",15,color,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def trainer_slide(kicker,name,role,rows,initials,accent=BLUE):
    """Profile-card layout: avatar badge + name/role panel on the left, labelled
    info tiles on the right. rows: list of (LABEL, value); blank value → fill-in line."""
    s=head(slide(),"About the Trainer",kicker,kcolor=accent)
    lx=Inches(0.85); lw=Inches(3.65)
    rect(s,lx,Inches(1.95),lw,Inches(4.7),_panel()); rect(s,lx,Inches(1.95),lw,Inches(0.12),accent)
    bd=Inches(1.7); ax=int(lx+(lw-bd)/2)
    oval(s,ax,Inches(2.5),bd,bd,accent)
    txt(s,ax,Inches(2.5),bd,bd,[[(initials,44,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,lx+Inches(0.15),Inches(4.55),lw-Inches(0.3),Inches(0.6),[[(name,21,_ink(),True)]],align=PP_ALIGN.CENTER)
    txt(s,lx+Inches(0.15),Inches(5.2),lw-Inches(0.3),Inches(1.2),[[(role,13,_grey(),False)]],align=PP_ALIGN.CENTER)
    rx=Inches(4.9); rw=Inches(7.6); ry=Inches(1.95); rh=Inches(4.7)
    n=len(rows); gy=Inches(0.2); th=int((rh-gy*(n-1))/n)
    for i,(label,val) in enumerate(rows):
        y=int(ry+(th+gy)*i); col=PALETTE[i%len(PALETTE)]
        rect(s,rx,y,rw,th,_panel()); rect(s,rx,y,Inches(0.1),th,col)
        vruns=[(val,14,_ink(),False)] if val else [("____________________________________________",13,_line(),False)]
        txt(s,rx+Inches(0.32),y,rw-Inches(0.6),th,
            [[(label.upper(),11,col,True)],vruns],anchor=MSO_ANCHOR.MIDDLE,space=3)
    footer(s); return s
def activity_overview(tag,title,desc,build,services,kicker):
    s=head(slide(),title,kicker,kcolor=TEAL)
    rect(s,Inches(0.85),Inches(1.85),Inches(1.7),Inches(0.5),_acc(TEAL))
    txt(s,Inches(0.85),Inches(1.9),Inches(1.7),Inches(0.4),[[(tag,16,DK_BG if THEME["dark"] else WHITE,True)]],align=PP_ALIGN.CENTER)
    # Auto-fit the description so long lab goals never clip into the panel below.
    dL=len(str(desc)); dsize = 21 if dL<=200 else (17 if dL<=320 else 15)
    txt(s,Inches(0.85),Inches(2.55),Inches(11.7),Inches(1.6),[[(desc,dsize,_ink(),False)]])
    rect(s,Inches(0.85),Inches(4.3),Inches(11.7),Inches(2.0),_panel())
    txt(s,Inches(1.1),Inches(4.5),Inches(11),Inches(0.4),[[("You'll build",14,_acc(BLUE),True)]])
    txt(s,Inches(1.1),Inches(4.9),Inches(11),Inches(0.6),[[(build,18,_ink(),True)]])
    txt(s,Inches(1.1),Inches(5.6),Inches(11.2),Inches(0.6),[[("Tools:  ",13,_grey(),True),(services,13,_grey(),False)]]); footer(s); return s
def step_slide(kicker,act_title,n,total,text,cmd=""):
    s=head(slide(),act_title,kicker,TEAL)
    oval(s,Inches(0.85),Inches(2.5),Inches(1.4),Inches(1.4),_acc(TEAL))
    txt(s,Inches(0.85),Inches(2.74),Inches(1.4),Inches(0.9),[[(str(n),38,DK_BG if THEME["dark"] else WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(0.95),Inches(1.95),Inches(11),Inches(0.4),[[(f"STEP {n} OF {total}",13,_grey(),True)]])
    txt(s,Inches(2.55),Inches(2.4),Inches(10.1),Inches(1.3),[[(text,23,_ink(),False)]],anchor=MSO_ANCHOR.MIDDLE)
    # HARD RULE: never render a comment-only "command" (# …) on a slide — the
    # code box appears only for a real, runnable command.
    if cmd and not cmd.lstrip().startswith("#"):
        if THEME["dark"]:
            # slightly lighter fill + faint border so the box reads against the dark bg
            rect(s,Inches(2.55),Inches(4.15),Inches(10.1),Inches(0.95),RGBColor(0x10,0x14,0x1E),line=DK_FAINT)
        else:
            rect(s,Inches(2.55),Inches(4.15),Inches(10.1),Inches(0.95),RGBColor(0x0B,0x12,0x20))
        txt(s,Inches(2.8),Inches(4.28),Inches(9.7),Inches(0.7),[[("$ "+cmd,13,RGBColor(0x9C,0xDC,0xFE),False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def dark_rows(tag,title_lines,sub_lines,rows,warn=None,accent=None,numbered=True):
    """Masterclass-style dark feature slide: terminal tag, amber (or custom
    accent) headline, light subtitle lines, numbered/bulleted row-cards. Rows
    may be strings, (amber,white) tuples, or lists of (text,colorkey) pairs
    (keys: w/a/b/g/d). warn=(BADGE, text) renders a red warning chip."""
    DKBG=RGBColor(0x0D,0x0E,0x0B); AMBER=RGBColor(0xF5,0xA6,0x23)
    ROWBG=RGBColor(0x15,0x17,0x12); FAINT=RGBColor(0x2A,0x2C,0x28); IVORY=RGBColor(0xE8,0xE6,0xDD)
    WARNRED=RGBColor(0xE0,0x5A,0x5A)
    ACC=accent or AMBER
    s=slide(); rect(s,0,0,SW,SH,DKBG)
    txt(s,Inches(0.85),Inches(0.4),Inches(11.6),Inches(0.35),[[("> "+tag,13,ACC,True)]])
    rect(s,Inches(0.85),Inches(0.82),Inches(11.63),Inches(0.02),FAINT)
    y=0.98
    for tl in title_lines:
        txt(s,Inches(0.85),Inches(y),Inches(11.9),Inches(0.78),[[(tl.upper(),40,ACC,True)]]); y+=0.76
    y+=0.08
    for sl in sub_lines:
        txt(s,Inches(0.85),Inches(y),Inches(11.9),Inches(0.38),[[(sl,15,IVORY,False)]]); y+=0.4
    y+=0.22
    n=max(1,len(rows)); gap=0.14
    bottom=6.35 if warn else 6.85
    rh=min(0.78,(bottom-y-gap*(n-1))/n)
    for i,r in enumerate(rows):
        ry=Inches(y+i*(rh+gap))
        rect(s,Inches(0.85),ry,Inches(11.63),Inches(rh),ROWBG)
        rect(s,Inches(0.85),ry,Inches(0.05),Inches(rh),ACC)
        if numbered:
            txt(s,Inches(1.1),ry,Inches(0.75),Inches(rh),[[(f"{i+1:02d}",16,ACC,True)]],anchor=MSO_ANCHOR.MIDDLE)
            tx,tw=Inches(1.9),Inches(10.35)
        else:
            txt(s,Inches(1.1),ry,Inches(0.4),Inches(rh),[[("▸",14,ACC,True)]],anchor=MSO_ANCHOR.MIDDLE)
            tx,tw=Inches(1.55),Inches(10.7)
        txt(s,tx,ry,tw,Inches(rh),[_dkruns(r,16)],anchor=MSO_ANCHOR.MIDDLE)
    if warn:
        wy=y+n*(rh+gap)+0.08
        rect(s,Inches(0.85),Inches(wy),Inches(1.15),Inches(0.34),WARNRED)
        txt(s,Inches(0.85),Inches(wy+0.02),Inches(1.15),Inches(0.3),[[(warn[0],12,DKBG,True)]],align=PP_ALIGN.CENTER)
        txt(s,Inches(2.2),Inches(wy),Inches(10.2),Inches(0.34),[[(warn[1],13,IVORY,False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
_DKC={"w":RGBColor(0xE8,0xE6,0xDD),"a":RGBColor(0xF5,0xA6,0x23),"b":RGBColor(0x6E,0xA8,0xFF),
      "g":RGBColor(0x4A,0xDE,0x80),"d":RGBColor(0x8A,0x8D,0x86),"r":RGBColor(0xE8,0x6A,0x8A)}
def _dkruns(spec,size=16):
    """spec: str | (amber,white) tuple | list of (text,colorkey) pairs -> txt() runs line."""
    if isinstance(spec,str): return [(spec,size,_DKC["w"],False)]
    if isinstance(spec,tuple): return [(spec[0],size,_DKC["a"],True),(spec[1],size,_DKC["w"],False)]
    return [(t,size,_DKC.get(k,_DKC["w"]),k in ("a","b","g","r")) for t,k in spec]
def dark_table(tag,title,col1,col2,rows,foot=None):
    """Masterclass-style dark comparison table: amber vs blue column headers,
    white row labels, thin divider rules, optional mixed-colour footer line."""
    DKBG=RGBColor(0x0D,0x0E,0x0B); AMBER=_DKC["a"]; HBLUE=_DKC["b"]
    FAINT=RGBColor(0x2A,0x2C,0x28); IVORY=_DKC["w"]; DIM=_DKC["d"]
    s=slide(); rect(s,0,0,SW,SH,DKBG)
    txt(s,Inches(0.85),Inches(0.4),Inches(8.0),Inches(0.35),[[("> "+tag,13,AMBER,True)]])
    txt(s,Inches(8.9),Inches(0.4),Inches(3.6),Inches(0.35),[[("● WHEN TO USE WHICH",12,DIM,True)]],align=PP_ALIGN.RIGHT)
    rect(s,Inches(0.85),Inches(0.82),Inches(11.63),Inches(0.02),FAINT)
    txt(s,Inches(0.85),Inches(1.0),Inches(11.9),Inches(0.8),[[(title.upper(),40,AMBER,True)]])
    LX,C1,C2=Inches(0.85),Inches(3.3),Inches(8.1)
    y=2.05
    txt(s,C1,Inches(y),Inches(4.6),Inches(0.4),[[(col1.upper(),15,AMBER,True)]])
    txt(s,C2,Inches(y),Inches(4.3),Inches(0.4),[[(col2.upper(),15,HBLUE,True)]])
    y+=0.5
    rh=0.62
    for label,v1,v2 in rows:
        rect(s,LX,Inches(y),Inches(11.63),Inches(0.015),FAINT)
        txt(s,LX,Inches(y+0.08),Inches(2.3),Inches(rh),[[(label,15,IVORY,True)]])
        txt(s,C1,Inches(y+0.08),Inches(4.6),Inches(rh),[[(v1,15,IVORY,False)]])
        txt(s,C2,Inches(y+0.08),Inches(4.3),Inches(rh),[[(v2,15,IVORY,False)]])
        y+=rh
    rect(s,LX,Inches(y),Inches(11.63),Inches(0.015),FAINT)
    if foot:
        txt(s,Inches(0.85),Inches(y+0.3),Inches(11.63),Inches(0.45),[_dkruns(foot,15)],align=PP_ALIGN.CENTER)
    footer(s); return s
def dark_code(tag,title,sub,code_lines,foot=None,win_title="",right_tag=None):
    """Masterclass-style dark slide with a terminal/code window: amber headline,
    window chrome (● ● ● + title), monospace coloured code lines, amber footer."""
    DKBG=RGBColor(0x0D,0x0E,0x0B); AMBER=_DKC["a"]; IVORY=_DKC["w"]; FAINT=RGBColor(0x2A,0x2C,0x28)
    CODEBG=RGBColor(0x08,0x0A,0x08); CHROME=RGBColor(0x14,0x16,0x12); DIM=_DKC["d"]
    s=slide(); rect(s,0,0,SW,SH,DKBG)
    txt(s,Inches(0.85),Inches(0.4),Inches(8.0),Inches(0.35),[[("> "+tag,13,AMBER,True)]])
    if right_tag:
        txt(s,Inches(8.9),Inches(0.4),Inches(3.6),Inches(0.35),[[("● "+right_tag,12,DIM,True)]],align=PP_ALIGN.RIGHT)
    rect(s,Inches(0.85),Inches(0.82),Inches(11.63),Inches(0.02),FAINT)
    txt(s,Inches(0.85),Inches(0.95),Inches(11.9),Inches(0.75),[[(title.upper(),38,AMBER,True)]])
    y=1.72
    if sub:
        txt(s,Inches(0.85),Inches(y),Inches(11.9),Inches(0.4),[[(sub,15,IVORY,False)]]); y+=0.46
    # window chrome
    ch=0.34
    rect(s,Inches(0.85),Inches(y),Inches(11.63),Inches(ch),CHROME)
    for k,cc in enumerate((RGBColor(0xE0,0x5A,0x5A),RGBColor(0xF5,0xA6,0x23),RGBColor(0x4A,0xDE,0x80))):
        oval(s,Inches(1.02+k*0.22),Inches(y+0.11),Inches(0.12),Inches(0.12),cc)
    if win_title:
        txt(s,Inches(1.8),Inches(y+0.02),Inches(9.5),Inches(0.3),[[(win_title,11,DIM,False)]])
    y+=ch
    n=len(code_lines); lh=min(0.34,(6.35-y)/max(1,n))
    rect(s,Inches(0.85),Inches(y),Inches(11.63),Inches(lh*n+0.2),CODEBG)
    for i,ln in enumerate(code_lines):
        runs=_dkruns(ln,13)
        # monospace look: Courier New via manual run construction
        tb=s.shapes.add_textbox(Inches(1.1),Inches(y+0.08+i*lh),Inches(11.1),Inches(lh))
        tf=tb.text_frame; tf.word_wrap=False
        p=tf.paragraphs[0]
        for t_,sz,col,bold in runs:
            r=p.add_run(); r.text=t_; r.font.size=Pt(12.5); r.font.bold=bold
            r.font.color.rgb=col; r.font.name="Courier New"
    if foot:
        txt(s,Inches(0.85),Inches(6.62),Inches(11.63),Inches(0.4),[[(foot,14,AMBER,True)]],align=PP_ALIGN.CENTER)
    footer(s); return s
def dark_catalog(tag,title,cards,scope=None,right_tag=None):
    """Masterclass-style dark slide: a row of amber-outlined cards, each with a
    heading + a list of tool names, plus an optional SCOPE note panel below."""
    DKBG=RGBColor(0x0D,0x0E,0x0B); AMBER=_DKC["a"]; IVORY=_DKC["w"]
    FAINT=RGBColor(0x2A,0x2C,0x28); CARDBG=RGBColor(0x10,0x12,0x0E); DIM=_DKC["d"]
    s=slide(); rect(s,0,0,SW,SH,DKBG)
    txt(s,Inches(0.85),Inches(0.4),Inches(8.0),Inches(0.35),[[("> "+tag,13,AMBER,True)]])
    if right_tag:
        txt(s,Inches(8.4),Inches(0.4),Inches(4.1),Inches(0.35),[[("● "+right_tag,12,DIM,True)]],align=PP_ALIGN.RIGHT)
    rect(s,Inches(0.85),Inches(0.82),Inches(11.63),Inches(0.02),FAINT)
    txt(s,Inches(0.85),Inches(0.98),Inches(11.9),Inches(0.78),[[(title.upper(),40,AMBER,True)]])
    y=1.95
    n=len(cards); gap=Inches(0.18); cw=int((Inches(11.63)-gap*(n-1))/n); chh=Inches(3.0); b=Inches(0.02)
    for i,(headt,items) in enumerate(cards):
        x=int(Inches(0.85)+i*(cw+gap)); yv=Inches(y)
        rect(s,x,yv,cw,chh,CARDBG)
        rect(s,x,yv,cw,b,AMBER); rect(s,x,int(yv+chh-b),cw,b,AMBER)
        rect(s,x,yv,b,chh,AMBER); rect(s,int(x+cw-b),yv,b,chh,AMBER)
        txt(s,x,int(yv+Inches(0.14)),cw,Inches(0.34),[[(headt.upper(),13,AMBER,True)]],align=PP_ALIGN.CENTER)
        rows=[[(it,12,IVORY if not it.startswith("+") and "·" not in it else DIM,False)] for it in items]
        txt(s,x+Inches(0.18),int(yv+Inches(0.58)),cw-Inches(0.3),int(chh-Inches(0.7)),rows,space=6)
    if scope:
        sy=y+3.0+0.25
        rect(s,Inches(0.85),Inches(sy),Inches(11.63),Inches(1.15),CARDBG)
        rect(s,Inches(0.85),Inches(sy),Inches(0.05),Inches(1.15),AMBER)
        txt(s,Inches(1.1),Inches(sy+0.12),Inches(1.3),Inches(0.35),[[("SCOPE",15,AMBER,True)]])
        txt(s,Inches(2.5),Inches(sy+0.1),Inches(9.7),Inches(0.98),[[(scope,12.5,IVORY,False)]])
    footer(s); return s
def dark_panels(tag,title,panels,foot=None,right_tag=None):
    """Masterclass-style dark slide: coloured outlined panels ('// HEAD' + body
    run-lines) and a mixed-colour footer. panels: (colorkey, head, [runlines])."""
    DKBG=RGBColor(0x0D,0x0E,0x0B); FAINT=RGBColor(0x2A,0x2C,0x28); PBG=RGBColor(0x10,0x12,0x0E)
    AMBER=_DKC["a"]; DIM=_DKC["d"]
    PCOL={"a":_DKC["a"],"b":_DKC["b"],"g":_DKC["g"],"r":RGBColor(0xE8,0x6A,0x8A)}
    s=slide(); rect(s,0,0,SW,SH,DKBG)
    txt(s,Inches(0.85),Inches(0.4),Inches(8.0),Inches(0.35),[[("> "+tag,13,AMBER,True)]])
    if right_tag:
        txt(s,Inches(8.4),Inches(0.4),Inches(4.1),Inches(0.35),[[("● "+right_tag,12,DIM,True)]],align=PP_ALIGN.RIGHT)
    rect(s,Inches(0.85),Inches(0.82),Inches(11.63),Inches(0.02),FAINT)
    txt(s,Inches(0.85),Inches(0.98),Inches(11.9),Inches(0.78),[[(title.upper(),40,AMBER,True)]])
    y=2.0; n=len(panels); gap=Inches(0.25)
    pw=int((Inches(11.63)-gap*(n-1))/n); ph=Inches(3.3); b=Inches(0.02)
    for i,(ck,head,body) in enumerate(panels):
        col=PCOL.get(ck,AMBER); x=int(Inches(0.85)+i*(pw+gap)); yv=Inches(y)
        rect(s,x,yv,pw,ph,PBG)
        rect(s,x,yv,pw,b,col); rect(s,x,int(yv+ph-b),pw,b,col)
        rect(s,x,yv,b,ph,col); rect(s,int(x+pw-b),yv,b,ph,col)
        txt(s,x+Inches(0.28),int(yv+Inches(0.2)),pw-Inches(0.5),Inches(0.4),[[("// "+head.upper(),15,col,True)]])
        txt(s,x+Inches(0.28),int(yv+Inches(0.72)),pw-Inches(0.52),int(ph-Inches(0.9)),
            [_dkruns(ln,13) for ln in body],space=6)
    if foot:
        txt(s,Inches(0.85),Inches(y+3.3+0.28),Inches(11.63),Inches(0.45),[_dkruns(foot,15)],align=PP_ALIGN.CENTER)
    footer(s); return s
def dark_flow(tag,title,steps,foot=None,right_tag=None):
    """Masterclass-style dark slide: amber-outlined step boxes joined by arrows,
    plus a mixed-colour footer. steps: (head, body)."""
    DKBG=RGBColor(0x0D,0x0E,0x0B); FAINT=RGBColor(0x2A,0x2C,0x28); PBG=RGBColor(0x10,0x12,0x0E)
    AMBER=_DKC["a"]; IVORY=_DKC["w"]; DIM=_DKC["d"]
    s=slide(); rect(s,0,0,SW,SH,DKBG)
    txt(s,Inches(0.85),Inches(0.4),Inches(8.0),Inches(0.35),[[("> "+tag,13,AMBER,True)]])
    if right_tag:
        txt(s,Inches(8.4),Inches(0.4),Inches(4.1),Inches(0.35),[[("● "+right_tag,12,DIM,True)]],align=PP_ALIGN.RIGHT)
    rect(s,Inches(0.85),Inches(0.82),Inches(11.63),Inches(0.02),FAINT)
    txt(s,Inches(0.85),Inches(1.05),Inches(11.9),Inches(0.85),[[(title.upper(),44,AMBER,True)]])
    y=2.55; n=len(steps); ag=Inches(0.38); b=Inches(0.02)
    bw=int((Inches(11.63)-ag*(n-1))/n); bh=Inches(1.55)
    for i,(head,body) in enumerate(steps):
        x=int(Inches(0.85)+i*(bw+ag)); yv=Inches(y)
        rect(s,x,yv,bw,bh,PBG)
        rect(s,x,yv,bw,b,AMBER); rect(s,x,int(yv+bh-b),bw,b,AMBER)
        rect(s,x,yv,b,bh,AMBER); rect(s,int(x+bw-b),yv,b,bh,AMBER)
        txt(s,x,int(yv+Inches(0.16)),bw,Inches(0.36),[[(head.upper(),14,AMBER,True)]],align=PP_ALIGN.CENTER)
        txt(s,x+Inches(0.12),int(yv+Inches(0.6)),bw-Inches(0.24),Inches(0.85),[[(body,11.5,IVORY,False)]],align=PP_ALIGN.CENTER)
        if i<n-1:
            txt(s,int(x+bw),int(yv+bh/2-Inches(0.2)),ag,Inches(0.4),[[("→",18,AMBER,True)]],align=PP_ALIGN.CENTER)
    if foot:
        txt(s,Inches(0.85),Inches(y+1.55+0.45),Inches(11.63),Inches(0.8),[_dkruns(foot,15)],align=PP_ALIGN.CENTER)
    footer(s); return s
def dark_cards(tag,title_lines,cards,notes=(),badge=None):
    """Masterclass-style dark slide: amber headline + a row of green-outlined cards
    (name + subtitle) + optional AUTO-INSTALLED badge line and dim footnotes."""
    DKBG=RGBColor(0x0D,0x0E,0x0B); AMBER=RGBColor(0xF5,0xA6,0x23); GREEN=RGBColor(0x4A,0xDE,0x80)
    IVORY=RGBColor(0xE8,0xE6,0xDD); FAINT=RGBColor(0x2A,0x2C,0x28); CARDBG=RGBColor(0x10,0x12,0x0E)
    DIM=RGBColor(0x8A,0x8D,0x86)
    s=slide(); rect(s,0,0,SW,SH,DKBG)
    txt(s,Inches(0.85),Inches(0.4),Inches(11.6),Inches(0.35),[[("> "+tag,13,AMBER,True)]])
    rect(s,Inches(0.85),Inches(0.82),Inches(11.63),Inches(0.02),FAINT)
    y=0.98
    for tl in title_lines:
        txt(s,Inches(0.85),Inches(y),Inches(11.9),Inches(0.78),[[(tl.upper(),40,AMBER,True)]]); y+=0.76
    y+=0.35
    n=len(cards); gap=Inches(0.25); cw=int((Inches(11.63)-gap*(n-1))/n); ch=Inches(1.55); b=Inches(0.025)
    for i,(name,sub) in enumerate(cards):
        x=int(Inches(0.85)+i*(cw+gap)); yv=Inches(y)
        rect(s,x,yv,cw,ch,CARDBG)
        rect(s,x,yv,cw,b,GREEN); rect(s,x,int(yv+ch-b),cw,b,GREEN)
        rect(s,x,yv,b,ch,GREEN); rect(s,int(x+cw-b),yv,b,ch,GREEN)
        txt(s,x,int(yv+Inches(0.2)),cw,Inches(0.62),[[(name,26,GREEN,True)]],align=PP_ALIGN.CENTER)
        txt(s,x,int(yv+Inches(0.92)),cw,Inches(0.42),[[(sub,15,IVORY,False)]],align=PP_ALIGN.CENTER)
    y+=1.55+0.4
    rest=list(notes)
    if badge and rest:
        rect(s,Inches(0.85),Inches(y),Inches(2.0),Inches(0.36),AMBER)
        txt(s,Inches(0.85),Inches(y+0.02),Inches(2.0),Inches(0.32),[[(badge,12,DKBG,True)]],align=PP_ALIGN.CENTER)
        txt(s,Inches(3.0),Inches(y),Inches(9.4),Inches(0.36),[[(rest.pop(0),14,IVORY,True)]],anchor=MSO_ANCHOR.MIDDLE)
        y+=0.5
    for nt in rest:
        txt(s,Inches(0.85),Inches(y),Inches(11.6),Inches(0.34),[[(nt,12,DIM,False)]]); y+=0.36
    footer(s); return s
def dark_pairs(tag,title_lines,sub,pairs):
    """Masterclass-style dark slide: amber headline + a 2-column grid of
    command/description cells (amber command, light description)."""
    import math as _m
    DKBG=RGBColor(0x0D,0x0E,0x0B); AMBER=RGBColor(0xF5,0xA6,0x23)
    IVORY=RGBColor(0xE8,0xE6,0xDD); FAINT=RGBColor(0x2A,0x2C,0x28); ROWBG=RGBColor(0x15,0x17,0x12)
    DIM=RGBColor(0x8A,0x8D,0x86)
    s=slide(); rect(s,0,0,SW,SH,DKBG)
    txt(s,Inches(0.85),Inches(0.4),Inches(11.6),Inches(0.35),[[("> "+tag,13,AMBER,True)]])
    rect(s,Inches(0.85),Inches(0.82),Inches(11.63),Inches(0.02),FAINT)
    y=0.98
    for tl in title_lines:
        txt(s,Inches(0.85),Inches(y),Inches(11.9),Inches(0.78),[[(tl.upper(),40,AMBER,True)]]); y+=0.76
    if sub:
        txt(s,Inches(0.85),Inches(y+0.02),Inches(11.9),Inches(0.36),[[(sub,14,DIM,False)]]); y+=0.44
    y+=0.18
    cols=2; rows=_m.ceil(len(pairs)/cols); gx=Inches(0.3); gy=0.13
    cw=int((Inches(11.63)-gx)/cols); rh=min(0.8,(6.8-y-gy*(rows-1))/rows)
    for i,(cmd,desc) in enumerate(pairs):
        r,c=i%rows,i//rows   # fill column-first like the reference
        x=int(Inches(0.85)+c*(cw+gx)); yv=Inches(y+r*(rh+gy))
        rect(s,x,yv,cw,Inches(rh),ROWBG)
        rect(s,x,yv,Inches(0.045),Inches(rh),AMBER)
        txt(s,x+Inches(0.22),yv,Inches(2.35),Inches(rh),[[(cmd,14,AMBER,True)]],anchor=MSO_ANCHOR.MIDDLE)
        txt(s,x+Inches(2.6),yv,cw-Inches(2.8),Inches(rh),[[(desc,12.5,IVORY,False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def shot(title,img,kicker=None,caption=""):
    """Screenshot slide — a framed, real UI capture with a caption (16:10 source)."""
    s=head(slide(),title,kicker,TEAL)
    H=Inches(4.5); W=int(H*1440/900)          # native 1440x900 aspect
    x=int((SW-W)/2); y=Inches(2.02)
    rect(s,x-Inches(0.06),y-Inches(0.06),W+Inches(0.12),H+Inches(0.12),_line())
    s.shapes.add_picture(img,x,y,height=H)
    if caption:
        txt(s,Inches(0.85),Inches(6.68),Inches(11.6),Inches(0.35),[[(caption,12,_grey(),False)]],align=PP_ALIGN.CENTER)
    footer(s); return s
def test_slide(act_title,text,kicker):
    s=head(slide(),act_title,kicker,TEAL)
    rect(s,Inches(0.85),Inches(2.3),Inches(11.7),Inches(2.6),DK_PANEL if THEME["dark"] else RGBColor(0xE8,0xF7,0xEE))
    txt(s,Inches(1.2),Inches(2.6),Inches(11),Inches(0.5),[[("✅  Test it",20,DK_GREEN if THEME["dark"] else RGBColor(0x12,0x7A,0x3E),True)]])
    txt(s,Inches(1.2),Inches(3.3),Inches(11),Inches(1.4),[[(text,18,_ink(),False)]]); footer(s); return s
def brk(kind,dur,color=AMBER):
    s=slide(); rect(s,0,0,SW,SH,_bg())
    rect(s,0,0,SW,Inches(0.22),color); rect(s,0,Inches(7.28),SW,Inches(0.22),color)
    rect(s,Inches(5.4),Inches(2.35),Inches(2.53),Inches(0.1),color)
    txt(s,0,Inches(2.75),SW,Inches(1.2),[[(kind,48,_ink(),True)]],align=PP_ALIGN.CENTER)
    txt(s,0,Inches(4.05),SW,Inches(0.8),[[(dur,22,color,True)]],align=PP_ALIGN.CENTER); PAGE["n"]+=1

# ============================================================ BUILD
cover()

# ---------------- ADMIN ----------------
section("COURSE ADMINISTRATION","Welcome & Housekeeping","")
trainer_slide("YOUR TRAINER · GENERAL","Your Trainer","General Trainer template —\nto be completed by the trainer",
 [("Name",""),("Title / Designation",""),("Qualifications",""),
  ("Areas of expertise",""),("Training & industry experience",""),("Contact","")],
 initials="?",accent=GREY)
trainer_slide("YOUR TRAINER",C.TRAINER,"Principal Trainer\nTertiary Infotech Academy Pte. Ltd.",
 [("Role","Principal Trainer, Tertiary Infotech Academy Pte. Ltd."),
  ("Certification",getattr(C,"TRAINER_CERT","Industry certified in the subject area of this course.")),
  ("Delivers",getattr(C,"TRAINER_DELIVERS","Professional short courses for Tertiary Infotech Academy.")),
  ("Founder","Founder and lead instructor at Tertiary Infotech / Tertiary Courses.")],
 initials="AA",accent=BLUE)
content("Let's Know Each Other",getattr(C,"ICE_BREAKER",[
 "Your name and organisation / role.",
 "Your experience with this subject (if any).",
 "What you want to be able to do after this course."]),kicker="ICE-BREAKER")
tile_grid("Ground Rules",[
 "Set your mobile phone to silent mode.","Participate actively — no question is too small.",
 "Mutual respect: agree to disagree.","One conversation at a time.",
 "Be punctual; return from breaks on time.","Learn by doing — try every hands-on lab."],
 kicker="HOUSEKEEPING",cols=2,size=15)
_dl=flow_h("Download Course Material",[
 "Sign in at the LMS-TMS portal",
 "Open this course from your dashboard",
 "Go to the Courseware tab",
 "Download the Slides (PPT/PDF), Learner Guide and Lesson Plan",
 "Keep them open alongside the labs as you work"],kicker="COURSE PORTAL")
# Full portal URL as a wide caption (kept out of the narrow chips so it never wraps mid-token)
txt(_dl,Inches(0.85),Inches(6.5),Inches(11.6),Inches(0.4),
    [[("Portal:  https://lms-tms.tertiaryinfotech.com",16,_acc(BLUE),True)]],align=PP_ALIGN.CENTER)
# Lesson plan overview — derived entirely from course_data (day themes + the
# topic→lab mapping) so it can never drift from the LP or the labs.
def _lab_range(acts):
    """Render a lab list compactly: contiguous runs collapse to 'Labs 2-5',
    gaps are listed ('Labs 2-5, 11'), and a single lab reads 'Lab 1'."""
    ns=sorted(a["num"] for a in acts)
    if not ns: return ""
    runs=[]; s=p=ns[0]
    for n in ns[1:]:
        if n==p+1: p=n; continue
        runs.append((s,p)); s=p=n
    runs.append((s,p))
    parts=[str(a) if a==b else f"{a}–{b}" for a,b in runs]
    return f" ({'Lab' if len(ns)==1 else 'Labs'} {', '.join(parts)})"
def _topic_line(t):
    return (f"Topic {t['num']}: {t['title']}{_lab_range(TOPIC_ACTS.get(t['num'],[]))}",1)
def _topics_by_day():
    """Which topics are taught on which day — read from the SAME schedule the
    Lesson Plan is built from, so the deck can never contradict the LP. Falls
    back to an even split only if the schedule can't be resolved."""
    try:
        def _lt(nums): return ""      # schedule text isn't needed, only the rows
        sched=C.SCHEDULE(_lt) if callable(getattr(C,"SCHEDULE",None)) else getattr(C,"SCHEDULE",{})
        by_day={}
        for day,(_theme,rows) in sched.items():
            found=[]
            for row in rows:
                if row[3] not in ("topic","lab"): continue
                text=row[4]; low=text.lower()
                for t in C.TOPICS:
                    if t["num"] in found: continue
                    head=t["title"].split(" — ")[0].strip()   # e.g. "Define" from "Define — Scope the Problem"
                    # Also match any significant word of the title as an
                    # ALL-CAPS phase label ("FOUNDATIONS", "MEASURE", …).
                    caps=[w.upper() for w in re.findall(r"[A-Za-z]{4,}",t["title"])]
                    if (f"topic {t['num']}" in low or t["title"].lower() in low
                            or head.upper() in text
                            or any(c in text for c in caps)):
                        found.append(t["num"])
            by_day[day]=found
        if any(by_day.values()):
            seen=set(); out={}
            for day in sorted(by_day):                 # a topic belongs to the first day it appears
                out[day]=[n for n in by_day[day] if not (n in seen or seen.add(n))]
            return out
    except Exception:
        pass
    half=(len(C.TOPICS)+1)//2
    return {1:[t["num"] for t in C.TOPICS[:half]], 2:[t["num"] for t in C.TOPICS[half:]]}

_BY_DAY=_topics_by_day()
_TB={t["num"]:t for t in C.TOPICS}
_d2=min(2,C.DAYS)
_left=[(f"Day 1 — {C.DAY_THEMES[1]}",0)]+[_topic_line(_TB[n]) for n in _BY_DAY.get(1,[]) if n in _TB]
_right=[(f"Day {_d2} — {C.DAY_THEMES[_d2]}",0)]+[_topic_line(_TB[n]) for n in _BY_DAY.get(_d2,[]) if n in _TB]+[
 ("Daily timing",0),("9:30am–6:30pm · 1-hour lunch · tea breaks within training time",1)]
two_col(f"Lesson Plan — {C.DAYS} Day{'s' if C.DAYS>1 else ''}, 8 hours/day",_left,_right,
 kicker="SCHEDULE",lhead="Day 1",rhead=f"Day {min(2,C.DAYS)}")
# Learning-outcome tiles built straight from course_data. Optional per-course
# short titles via course_data.LO_TITLES; otherwise fall back to "LO1", "LO2", …
_LO_TITLES=getattr(C,"LO_TITLES",[])
_lo_tiles=[(_LO_TITLES[i] if i<len(_LO_TITLES) else f"LO{i+1}",
            lo.split(": ",1)[1] if ": " in lo else lo)
           for i,lo in enumerate(C.LEARNING_OUTCOMES)]
tile_grid("Learning Outcomes",_lo_tiles,kicker="WHAT YOU'LL ACHIEVE",cols=2,size=15)
# NON-WSQ: there is NO assessment. Instead, set out how learning is reinforced
# in class — practice, feedback and self-check against the learning outcomes.
flow_h("How You'll Learn",[
 "Trainer demonstrates the concept",
 "You build it yourself in the hands-on lab",
 "Verify with the lab's 'Test it' step",
 "Discuss and troubleshoot together",
 "Recap against the learning outcomes"],kicker="LEARNING APPROACH")

# ---------------- CORE CONCEPTS ----------------
# Generic, data-driven course-concepts section. A course supplies its own
# framing via optional course_data fields; nothing here is tied to a specific
# subject, so the same engine builds any non-WSQ course.
_OVERVIEW=getattr(C,"COURSE_OVERVIEW",None)
if _OVERVIEW:
    section("CORE CONCEPTS",_OVERVIEW.get("section_title","Course Fundamentals"),"")
    if _OVERVIEW.get("concepts"):
        tile_grid(_OVERVIEW.get("concepts_title","Key Concepts"),_OVERVIEW["concepts"],
                  kicker="OVERVIEW",cols=2,size=15)
    if _OVERVIEW.get("framework"):
        tile_grid(_OVERVIEW.get("framework_title","The Framework"),_OVERVIEW["framework"],
                  kicker="THE BUILDING BLOCKS",cols=2,size=14)
    if _OVERVIEW.get("statement"):
        st=_OVERVIEW["statement"]
        big_statement(st.get("headline",""),st.get("body",""),st.get("kicker",""),color=BLUE)
    if _OVERVIEW.get("pillars"):
        cards3(_OVERVIEW.get("pillars_title","What You'll Build"),
               [(CARD_COLORS[i%3],t,b) for i,(t,b) in enumerate(_OVERVIEW["pillars"])],
               kicker="COURSE ROADMAP")
    if _OVERVIEW.get("arc"):
        content(_OVERVIEW.get("arc_title","How Every Lab Progresses"),_OVERVIEW["arc"],
                kicker="THE LEARNING ARC")

# ---------------- TOPICS + ACTIVITIES ----------------
# Optional per-lab screenshot slides. A course drops images into
# courseware/assets/screenshots/ and maps them in course_data.LAB_SHOTS:
#   LAB_SHOTS = {2: [("file.png", "Label", "Caption"), ...]}
_SHOTDIR=os.path.join(REPO,"courseware","assets","screenshots")
_SHOTS=getattr(C,"LAB_SHOTS",{})
def _lab_extras(num):
    for fn,label,cap in _SHOTS.get(num,[]):
        p=os.path.join(_SHOTDIR,fn)
        if os.path.exists(p):
            shot(f"{label}",p,kicker=f"LAB {num} · SCREENSHOT",caption=cap)


for t in C.TOPICS:
    accent=TOPIC_THEME.get(t["num"],BLUE)
    section(f"TOPIC {t['code']}", t["title"], t["code"], t["subtitle"], accent=accent)
    # concept slide(s) — visual tile grid instead of a bullet list (topic-accented)
    # Non-WSQ: no exam, so the kicker shows the topic's share of course time
    # (course_data may omit `weighting` entirely).
    _w=t.get("weighting")
    tile_grid(f"Key Concepts — {t['title']}", t["concepts"],
              kicker=(f"COURSE COVERAGE {_w}" if _w else "KEY CONCEPTS"),
              cols=2, size=14, accent=accent)
    acts=TOPIC_ACTS[t["num"]]
    # a card summary of the labs in this topic
    # Split the topic's labs into at most three cards. Only emit cards that
    # actually hold labs — a topic with one or two labs renders one or two
    # cards rather than padding the row with empty "—" placeholders.
    third=max(1,(len(acts)+2)//3)
    groups=[g for g in (acts[i:i+third] for i in range(0,len(acts),third)) if g][:3]
    cards=[(CARD_COLORS[gi], _lab_range(g).strip(" ()"), [a["title"] for a in g])
           for gi,g in enumerate(groups)]
    cards3(f"Hands-On Labs — {t['title']}", cards, kicker="WHAT YOU'LL DO")
    # per activity
    for a in acts:
        activity_overview(f"LAB {a['num']}", a["title"], a["desc"], a["build"], a["services"], kicker=f"TOPIC {t['code']} · HANDS-ON")
        _lab_extras(a["num"])
        # The deck never shows YouTube reference links — videos live in the labs/ READMEs only.
        steps=[(si,sc) for (si,sc) in a["steps"] if "youtube" not in sc.lower() and "youtube" not in si.lower()]
        total=len(steps)
        for i,(instr,cmd) in enumerate(steps,1):
            step_slide(f"LAB {a['num']}", a["title"], i, total, instr, cmd)
        test_slide(a["title"], a["test"], kicker=f"LAB {a['num']} · VERIFY")
    # topic recap
    content(f"Recap — {t['title']}",
            ["You can now: "+a["objective"] for a in {x["objective"]:x for x in acts}.values()][:6],
            kicker="TOPIC RECAP", size=17)

# ---------------- CLOSE ----------------
section("WRAP-UP","Course Summary & Next Steps","")
# "What You Achieved" tiles built straight from the learning outcomes so they can never drift
# Achievement tiles reuse the learning-outcome titles so they can never drift.
_ach_tiles=[(_LO_TITLES[i] if i<len(_LO_TITLES) else f"LO{i+1}",
             lo.split(": ",1)[1] if ": " in lo else lo)
            for i,lo in enumerate(C.LEARNING_OUTCOMES)]
tile_grid("What You Achieved",_ach_tiles,kicker="LEARNING OUTCOMES",cols=2,size=15)
_next=getattr(C,"NEXT_STEPS",None)
if _next:
    content(_next.get("title","Continuing Your Journey"),_next["items"],kicker="NEXT STEPS")
content("Keep Practising After the Course",[
 "Re-run the labs on your own machine — repetition is what makes the skills stick.",
 "Apply the techniques to a real process or project in your own organisation.",
 "Your course material stays available at https://lms-tms.tertiaryinfotech.com/.",
 "Explore the follow-on courses at www.tertiarycourses.com.sg."],kicker="WRAP-UP")
_ty=getattr(C,"THANK_YOU",{})
big_statement("Thank You!",
              _ty.get("body","You now have the practical skills from every lab in this course — keep applying them to your own work."),
              _ty.get("kicker","THANK YOU FOR ATTENDING"),color=TEAL)

# Fade slide transition on every slide — matches the source masterclass deck's
# fade aesthetic. python-pptx has no transition API, so inject the XML directly.
from lxml import etree as _et
from pptx.oxml.ns import qn as _qn
_TRANS=('<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        'spd="med"><p:fade/></p:transition>')
for _s in prs.slides:
    if _s.element.find(_qn('p:transition')) is None:
        _s.element.append(_et.fromstring(_TRANS))

OUT=os.path.join(REPO,"courseware",f"{C.SHORT_TITLE}-{C.VERSION}.pptx")
prs.save(OUT)
print(f"Saved {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
